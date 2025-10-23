import json
import os
import tkinter
from tkinter import filedialog, messagebox
import customtkinter as ctk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


def criar_apresentacao_minimalista(caminho_json, caminho_imagem, font_size_title, font_size_artist, font_size_lyrics):
    """
    Cria uma apresentação minimalista a partir de um JSON,
    com tamanhos de fonte personalizáveis.
    """
    try:
        with open(caminho_json, 'r', encoding='utf-8') as f:
            dados_musica = json.load(f)

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        cor_texto = RGBColor(255, 255, 255)
        left = top = Inches(0)

    
        slide_titulo_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_titulo_layout)
        slide.shapes.add_picture(caminho_imagem, left, top, width=prs.slide_width, height=prs.slide_height)


        titulo_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
        p_titulo = titulo_shape.text_frame.paragraphs[0]
        p_titulo.text = dados_musica.get('title', 'Título Desconhecido')
        p_titulo.font.bold = True
        p_titulo.font.size = Pt(font_size_title) 
        p_titulo.font.color.rgb = cor_texto
        p_titulo.alignment = PP_ALIGN.CENTER

  
        subtitulo_shape = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(14), Inches(2))
        p_subtitulo = subtitulo_shape.text_frame.paragraphs[0]
        p_subtitulo.text = dados_musica.get('artist', '')
        p_subtitulo.font.size = Pt(font_size_artist) 
        p_subtitulo.font.color.rgb = cor_texto
        p_subtitulo.alignment = PP_ALIGN.CENTER

     
        slide_letra_layout = prs.slide_layouts[6]
        paragraphs = dados_musica.get('lyrics', {}).get('paragraphs', [])
        
        for paragrafo in paragraphs:
            slide_letra = prs.slides.add_slide(slide_letra_layout)
            slide_letra.shapes.add_picture(caminho_imagem, left, top, width=prs.slide_width, height=prs.slide_height)

            caixa_texto = slide_letra.shapes.add_textbox(Inches(1), Inches(0), Inches(14), prs.slide_height)
            frame_texto = caixa_texto.text_frame
            frame_texto.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            frame_texto.word_wrap = True
            
            p = frame_texto.paragraphs[0]
            p.text = paragrafo.get('text', '')
            p.font.size = Pt(font_size_lyrics) 
            p.font.bold = True
            p.font.color.rgb = cor_texto
            p.alignment = PP_ALIGN.CENTER

   
        diretorio, nome_arquivo = os.path.split(caminho_json)
        nome_base = os.path.splitext(nome_arquivo)[0]
        caminho_saida = os.path.join(diretorio, f"{nome_base}_Minimalista.pptx")
        prs.save(caminho_saida)
        
        return True, caminho_saida 

    except Exception as e:
        return False, f"Erro ao processar {os.path.basename(caminho_json)}: {e}" 


class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("Gerador de Apresentações")
        self.geometry("600x450")
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("blue")

        self.json_files = []
        self.image_path = ""

        self.grid_columnconfigure(1, weight=1)



        self.btn_select_jsons = ctk.CTkButton(self, text="1. Selecionar Arquivos JSON", command=self.select_json_files)
        self.btn_select_jsons.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        
        self.lbl_json_status = ctk.CTkLabel(self, text="Nenhum arquivo selecionado", anchor="w")
        self.lbl_json_status.grid(row=0, column=1, padx=20, pady=(20, 10), sticky="ew")

        self.btn_select_image = ctk.CTkButton(self, text="2. Selecionar Imagem de Fundo", command=self.select_image_file)
        self.btn_select_image.grid(row=1, column=0, padx=20, pady=10, sticky="ew")

        self.lbl_image_status = ctk.CTkLabel(self, text="Nenhuma imagem selecionada", anchor="w")
        self.lbl_image_status.grid(row=1, column=1, padx=20, pady=10, sticky="ew")


        self.lbl_font_title = ctk.CTkLabel(self, text="Tamanho Fonte (Título):")
        self.lbl_font_title.grid(row=2, column=0, padx=20, pady=10, sticky="w")
        self.entry_title = ctk.CTkEntry(self, placeholder_text="Ex: 88")
        self.entry_title.insert(0, "88")
        self.entry_title.grid(row=2, column=1, padx=20, pady=10, sticky="ew")

        self.lbl_font_artist = ctk.CTkLabel(self, text="Tamanho Fonte (Artista):")
        self.lbl_font_artist.grid(row=3, column=0, padx=20, pady=10, sticky="w")
        self.entry_artist = ctk.CTkEntry(self, placeholder_text="Ex: 40")
        self.entry_artist.insert(0, "40")
        self.entry_artist.grid(row=3, column=1, padx=20, pady=10, sticky="ew")

        self.lbl_font_lyrics = ctk.CTkLabel(self, text="Tamanho Fonte (Letra):")
        self.lbl_font_lyrics.grid(row=4, column=0, padx=20, pady=10, sticky="w")
        self.entry_lyrics = ctk.CTkEntry(self, placeholder_text="Ex: 96")
        self.entry_lyrics.insert(0, "96")
        self.entry_lyrics.grid(row=4, column=1, padx=20, pady=10, sticky="ew")

        self.generate_button = ctk.CTkButton(self, text="✨ Gerar Apresentações", command=self.start_generation, height=40)
        self.generate_button.grid(row=5, column=0, columnspan=2, padx=20, pady=20, sticky="ew")

        self.status_label = ctk.CTkLabel(self, text="")
        self.status_label.grid(row=6, column=0, columnspan=2, padx=20, pady=10, sticky="ew")

    def select_json_files(self):

        self.json_files = filedialog.askopenfilenames(
            title="Selecione os arquivos JSON das músicas",
            filetypes=[("Arquivos JSON", "*.json")]
        )
        if self.json_files:
            self.lbl_json_status.configure(text=f"{len(self.json_files)} arquivos selecionados")
        else:
            self.lbl_json_status.configure(text="Nenhum arquivo selecionado")

    def select_image_file(self):
        self.image_path = filedialog.askopenfilename(
            title="Selecione a imagem de fundo",
            filetypes=[("Imagens", "*.jpg *.jpeg *.png")]
        )
        if self.image_path:
            self.lbl_image_status.configure(text=os.path.basename(self.image_path))
        else:
            self.lbl_image_status.configure(text="Nenhuma imagem selecionada")

    def start_generation(self):
        if not self.json_files:
            messagebox.showwarning("Atenção", "Por favor, selecione pelo menos um arquivo JSON.")
            return
        if not self.image_path:
            messagebox.showwarning("Atenção", "Por favor, selecione uma imagem de fundo.")
            return

        try:
            font_title = int(self.entry_title.get())
            font_artist = int(self.entry_artist.get())
            font_lyrics = int(self.entry_lyrics.get())
        except ValueError:
            messagebox.showerror("Erro", "Os tamanhos de fonte devem ser números inteiros.")
            return


        self.generate_button.configure(state="disabled", text="Gerando...")
        success_count = 0
        errors = []

        total_files = len(self.json_files)
        for i, json_path in enumerate(self.json_files):
            self.status_label.configure(text=f"Processando {i+1}/{total_files}: {os.path.basename(json_path)}")
            self.update_idletasks() 
            
            success, message = criar_apresentacao_minimalista(
                json_path, self.image_path, font_title, font_artist, font_lyrics
            )
            if success:
                success_count += 1
            else:
                errors.append(message)

        self.generate_button.configure(state="normal", text="Gerar Apresentações")
        self.status_label.configure(text="")
        
        final_message = f"{success_count} de {total_files} apresentações geradas com sucesso!"
        if errors:
            final_message += "\n\nOcorreram os seguintes erros:\n- " + "\n- ".join(errors)
            messagebox.showerror("Processo Concluído com Erros", final_message)
        else:
            messagebox.showinfo("Processo Concluído", final_message)


if __name__ == "__main__":
    app = App()
    app.mainloop()